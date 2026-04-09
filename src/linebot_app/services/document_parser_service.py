from __future__ import annotations

from io import BytesIO
from pathlib import Path


class DocumentParserService:
    """解析常見文件檔案並抽取純文字。"""

    def __init__(self, *, max_chars: int = 6000) -> None:
        self.max_chars = max_chars

    def extract_text(self, *, file_name: str, content: bytes) -> tuple[str, str | None]:
        """
        Returns:
            (text, error)
            - text: 解析後文字（失敗時為空）
            - error: 失敗原因（成功時為 None）
        """
        suffix = Path(file_name or "").suffix.lower()

        try:
            if suffix == ".pdf":
                text = self._parse_pdf(content)
            elif suffix == ".docx":
                text = self._parse_docx(content)
            elif suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
                text = self._parse_excel(content)
            elif suffix == ".pptx":
                text = self._parse_pptx(content)
            elif suffix in {".txt", ".md", ".csv", ".tsv"}:
                text = self._parse_text(content)
            elif suffix in {".doc", ".xls", ".ppt"}:
                return (
                    "",
                    "目前不支援舊版 Office 二進位格式（.doc/.xls/.ppt），"
                    "請轉成 .docx/.xlsx/.pptx。",
                )
            else:
                return "", f"目前不支援此檔案格式：{suffix or '未知格式'}"
        except ModuleNotFoundError:
            return "", "檔案解析套件尚未安裝完成，請先執行 uv sync。"
        except Exception as exc:
            return "", f"檔案解析失敗：{exc}"

        text = (text or "").strip()
        if not text:
            return "", "檔案內容讀取成功，但未辨識到可用文字。"
        if len(text) > self.max_chars:
            return text[: self.max_chars], None
        return text, None

    def _parse_pdf(self, content: bytes) -> str:
        from pypdf import PdfReader  # type: ignore[import-untyped]

        reader = PdfReader(BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)

    def _parse_docx(self, content: bytes) -> str:
        from docx import Document  # type: ignore[import-untyped]

        doc = Document(BytesIO(content))
        lines = [p.text for p in doc.paragraphs if (p.text or "").strip()]
        return "\n".join(lines)

    def _parse_excel(self, content: bytes) -> str:
        from openpyxl import load_workbook  # type: ignore[import-untyped]

        wb = load_workbook(filename=BytesIO(content), read_only=True, data_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            lines.append(f"[工作表] {ws.title}")
            for row in ws.iter_rows(values_only=True):
                values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if values:
                    lines.append("\t".join(values))
        return "\n".join(lines)

    def _parse_pptx(self, content: bytes) -> str:
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(BytesIO(content))
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"[投影片 {i}]")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())
        return "\n".join(lines)

    def _parse_text(self, content: bytes) -> str:
        for encoding in ("utf-8", "utf-8-sig", "cp950", "big5", "latin-1"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="ignore")
